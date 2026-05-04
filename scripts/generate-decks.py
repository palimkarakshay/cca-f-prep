#!/usr/bin/env python3
"""
Generate .pptx and self-contained .html (reveal.js) versions of each
markdown deck under plans/. Re-run after editing the markdown source
to regenerate the artifacts.

Usage:  python3 scripts/generate-decks.py
"""

from __future__ import annotations

import json
import re
from html import escape as html_escape
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
PLANS = ROOT / "plans"

DECKS = [
    "deck-overview.md",
    "deck-investor.md",
    "deck-b2b-prospect.md",
    "deck-collaborator.md",
]


def split_slides(md: str) -> list[str]:
    """Split a markdown deck on horizontal rule slide separators.

    A separator is a line of exactly '---' (with optional whitespace),
    that is not inside a fenced code block.
    """
    slides: list[str] = []
    current: list[str] = []
    in_fence = False

    for raw in md.splitlines():
        line = raw.rstrip("\n")
        if line.strip().startswith("```"):
            in_fence = not in_fence
            current.append(line)
            continue
        if not in_fence and re.fullmatch(r"\s*---\s*", line):
            slides.append("\n".join(current).strip())
            current = []
            continue
        current.append(line)

    if current:
        slides.append("\n".join(current).strip())

    return [s for s in slides if s]


def strip_leading_blockquote(md: str) -> str:
    """Drop the first contiguous blockquote (the companion-docs note)."""
    lines = md.splitlines()
    out: list[str] = []
    i = 0
    n = len(lines)
    # skip leading blank lines
    while i < n and not lines[i].strip():
        i += 1
    # skip the first blockquote block
    if i < n and lines[i].lstrip().startswith(">"):
        while i < n and (lines[i].lstrip().startswith(">") or not lines[i].strip()):
            i += 1
    out.extend(lines[i:])
    return "\n".join(out)


# ---------------------------------------------------------------------------
# .pptx generation
# ---------------------------------------------------------------------------

def parse_slide_blocks(slide_md: str) -> tuple[str | None, list[dict]]:
    """Return (title, blocks). Each block is one of:
    - {"kind": "para", "text": str}
    - {"kind": "bullets", "items": [(level, text), ...]}
    - {"kind": "table", "rows": [[str, ...], ...]}
    - {"kind": "code", "text": str}
    - {"kind": "quote", "text": str}
    """
    title: str | None = None
    blocks: list[dict] = []

    lines = slide_md.splitlines()
    i = 0
    n = len(lines)

    while i < n:
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # Heading -> title (first one wins)
        m = re.match(r"^(#{1,6})\s+(.+?)\s*$", line)
        if m:
            if title is None:
                title = m.group(2)
            else:
                blocks.append({"kind": "para", "text": m.group(2), "bold": True})
            i += 1
            continue

        # Fenced code
        if stripped.startswith("```"):
            i += 1
            buf: list[str] = []
            while i < n and not lines[i].strip().startswith("```"):
                buf.append(lines[i])
                i += 1
            if i < n:
                i += 1  # consume closing fence
            blocks.append({"kind": "code", "text": "\n".join(buf)})
            continue

        # Blockquote
        if stripped.startswith(">"):
            buf = []
            while i < n and lines[i].strip().startswith(">"):
                buf.append(re.sub(r"^\s*>\s?", "", lines[i]))
                i += 1
            blocks.append({"kind": "quote", "text": "\n".join(buf).strip()})
            continue

        # Table (simple: a line that contains '|' and the next line is a separator)
        if "|" in stripped and i + 1 < n and re.match(r"^\s*\|?\s*[:\- |]+\|?\s*$", lines[i + 1]):
            header = parse_table_row(lines[i])
            i += 2  # skip separator
            rows = [header]
            while i < n and "|" in lines[i] and lines[i].strip():
                rows.append(parse_table_row(lines[i]))
                i += 1
            blocks.append({"kind": "table", "rows": rows})
            continue

        # Bullet list (collect contiguous bullets, supporting two indent levels)
        if re.match(r"^(\s*)([-*])\s+", line) or re.match(r"^(\s*)\d+\.\s+", line):
            items: list[tuple[int, str]] = []
            while i < n:
                bm = re.match(r"^(\s*)([-*]|\d+\.)\s+(.*)$", lines[i])
                if not bm:
                    if not lines[i].strip():
                        # allow a blank line inside the list, then peek
                        if i + 1 < n and re.match(r"^(\s*)([-*]|\d+\.)\s+", lines[i + 1]):
                            i += 1
                            continue
                    break
                indent = len(bm.group(1).expandtabs(4))
                level = 0 if indent < 2 else 1
                items.append((level, bm.group(3)))
                i += 1
            blocks.append({"kind": "bullets", "items": items})
            continue

        # Plain paragraph (one line; collapse soft-wrapped lines)
        para = [line]
        i += 1
        while i < n and lines[i].strip() and not re.match(r"^(#{1,6}\s|>|```|[-*]\s|\d+\.\s)", lines[i].lstrip()) and "|" not in lines[i]:
            para.append(lines[i])
            i += 1
        blocks.append({"kind": "para", "text": " ".join(s.strip() for s in para)})

    return title, blocks


def parse_table_row(line: str) -> list[str]:
    parts = line.strip().strip("|").split("|")
    return [p.strip() for p in parts]


# Minimal inline-markdown stripper for plaintext rendering in pptx
INLINE_BOLD = re.compile(r"\*\*(.+?)\*\*")
INLINE_ITALIC = re.compile(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)")
INLINE_CODE = re.compile(r"`([^`]+)`")
INLINE_LINK = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")

def md_to_text(s: str) -> str:
    s = INLINE_LINK.sub(r"\1", s)
    s = INLINE_BOLD.sub(r"\1", s)
    s = INLINE_ITALIC.sub(r"\1", s)
    s = INLINE_CODE.sub(r"\1", s)
    return s


def md_to_runs(p, s: str) -> None:
    """Add runs to a pptx paragraph, preserving bold/italic/code styling."""
    pos = 0
    pattern = re.compile(
        r"(\*\*(?P<bold>[^*]+)\*\*)"
        r"|(\*(?P<italic>[^*]+)\*)"
        r"|(`(?P<code>[^`]+)`)"
        r"|(\[(?P<linktext>[^\]]+)\]\([^)]+\))"
    )
    for m in pattern.finditer(s):
        if m.start() > pos:
            run = p.add_run()
            run.text = s[pos:m.start()]
        group = m.lastgroup
        text = m.group(group)
        run = p.add_run()
        run.text = text
        if group == "bold":
            run.font.bold = True
        elif group == "italic":
            run.font.italic = True
        elif group == "code":
            run.font.name = "Consolas"
        pos = m.end()
    if pos < len(s):
        run = p.add_run()
        run.text = s[pos:]


def build_pptx(deck_path: Path, out_path: Path) -> None:
    md = strip_leading_blockquote(deck_path.read_text(encoding="utf-8"))
    slides = split_slides(md)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for idx, slide_md in enumerate(slides):
        title, blocks = parse_slide_blocks(slide_md)
        slide = prs.slides.add_slide(blank_layout)

        # Title
        if title:
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
            tf = tx.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            md_to_runs(p, md_to_text(title))
            for run in p.runs:
                run.font.size = Pt(34 if idx == 0 else 28)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        # Body
        body_top = Inches(1.4 if idx > 0 else 1.6)
        body_box = slide.shapes.add_textbox(
            Inches(0.5), body_top, Inches(12.3), Inches(7.5) - body_top - Inches(0.3)
        )
        body_tf = body_box.text_frame
        body_tf.word_wrap = True

        first_para_used = False

        def new_p():
            nonlocal first_para_used
            if not first_para_used:
                p = body_tf.paragraphs[0]
                first_para_used = True
            else:
                p = body_tf.add_paragraph()
            return p

        for block in blocks:
            if block["kind"] == "para":
                p = new_p()
                p.alignment = PP_ALIGN.LEFT
                md_to_runs(p, md_to_text(block["text"]))
                for run in p.runs:
                    run.font.size = Pt(18)
                    if block.get("bold"):
                        run.font.bold = True
                p.space_after = Pt(8)

            elif block["kind"] == "bullets":
                for level, item in block["items"]:
                    p = new_p()
                    p.level = level
                    md_to_runs(p, md_to_text(item))
                    for run in p.runs:
                        run.font.size = Pt(16 if level == 0 else 14)
                p.space_after = Pt(6)

            elif block["kind"] == "code":
                p = new_p()
                run = p.add_run()
                run.text = block["text"]
                run.font.name = "Consolas"
                run.font.size = Pt(11)
                p.space_after = Pt(6)

            elif block["kind"] == "quote":
                p = new_p()
                md_to_runs(p, md_to_text(block["text"]))
                for run in p.runs:
                    run.font.size = Pt(18)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(0x4A, 0x4A, 0x6E)
                p.space_after = Pt(6)

            elif block["kind"] == "table":
                rows = block["rows"]
                if not rows:
                    continue
                cols = max(len(r) for r in rows)
                # add table after current text content, sizing roughly
                left = Inches(0.5)
                top = body_top + Inches(0.6 + 0.45 * sum(
                    (1 if b["kind"] == "para" else
                     len(b.get("items", [])) if b["kind"] == "bullets" else 1)
                    for b in blocks[: blocks.index(block)]
                ))
                width = Inches(12.3)
                height = Inches(min(0.45 * len(rows), 5.0))
                tbl_shape = slide.shapes.add_table(len(rows), cols, left, top, width, height)
                tbl = tbl_shape.table
                for r_idx, row in enumerate(rows):
                    for c_idx in range(cols):
                        cell = tbl.cell(r_idx, c_idx)
                        cell_text = row[c_idx] if c_idx < len(row) else ""
                        tf = cell.text_frame
                        tf.text = ""
                        cp = tf.paragraphs[0]
                        md_to_runs(cp, md_to_text(cell_text))
                        for run in cp.runs:
                            run.font.size = Pt(11)
                            if r_idx == 0:
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        if r_idx == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    prs.save(out_path)


# ---------------------------------------------------------------------------
# Self-contained reveal.js HTML generation
# ---------------------------------------------------------------------------

REVEAL_HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>__TITLE__</title>
  <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.1/reveal.min.css" />
  <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.1/theme/white.min.css" id="theme" />
  <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.10.0/styles/github.min.css" />
  <style>
    .reveal { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif; }
    .reveal h1, .reveal h2, .reveal h3 { color: #1a1a2e; text-transform: none; letter-spacing: 0; font-weight: 700; }
    .reveal h1 { font-size: 2.4em; }
    .reveal h2 { font-size: 1.7em; }
    .reveal h3 { font-size: 1.2em; }
    .reveal section { text-align: left; font-size: 0.7em; }
    .reveal table { border-collapse: collapse; margin: 0.6em 0; font-size: 0.85em; width: 100%; }
    .reveal table th, .reveal table td { border: 1px solid #ccc; padding: 6px 10px; vertical-align: top; }
    .reveal table th { background: #1a1a2e; color: #fff; }
    .reveal pre { font-size: 0.6em; }
    .reveal code { font-family: "SF Mono", Menlo, Consolas, monospace; }
    .reveal blockquote { border-left: 4px solid #1a1a2e; background: #f4f4f8; padding: 0.6em 1em; font-style: italic; }
    .reveal ul, .reveal ol { margin-left: 1.2em; }
  </style>
</head>
<body>
  <div class="reveal">
    <div class="slides">
      <section data-markdown data-separator="^---$">
        <textarea data-template>
__MD__
        </textarea>
      </section>
    </div>
  </div>
  <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.1/reveal.min.js"></script>
  <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.1/plugin/markdown/markdown.min.js"></script>
  <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.1/plugin/highlight/highlight.min.js"></script>
  <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.6.1/plugin/notes/notes.min.js"></script>
  <script>
    Reveal.initialize({
      hash: true,
      slideNumber: 'c/t',
      transition: 'fade',
      width: 1280,
      height: 800,
      margin: 0.06,
      plugins: [ RevealMarkdown, RevealHighlight, RevealNotes ]
    });
  </script>
</body>
</html>
"""


def build_html(deck_path: Path, out_path: Path) -> None:
    md = strip_leading_blockquote(deck_path.read_text(encoding="utf-8"))
    # Title from first H1
    title_match = re.search(r"^#\s+(.+)$", md, re.MULTILINE)
    title = title_match.group(1).strip() if title_match else deck_path.stem

    # The reveal.js markdown plugin reads the textarea raw; we just need to
    # avoid the textarea closing tag appearing inside the content.
    safe_md = md.replace("</textarea>", "<\\/textarea>")

    html = (REVEAL_HTML_TEMPLATE
            .replace("__TITLE__", html_escape(title))
            .replace("__MD__", safe_md))
    out_path.write_text(html, encoding="utf-8")


# ---------------------------------------------------------------------------

def main() -> None:
    for fname in DECKS:
        src = PLANS / fname
        if not src.exists():
            print(f"skip: {fname} not found")
            continue
        stem = src.stem  # e.g. "deck-overview"
        pptx_out = PLANS / f"{stem}.pptx"
        html_out = PLANS / f"{stem}.html"
        build_pptx(src, pptx_out)
        build_html(src, html_out)
        print(f"generated: {pptx_out.name}, {html_out.name}")


if __name__ == "__main__":
    main()
